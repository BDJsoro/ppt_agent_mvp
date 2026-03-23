import os
import datetime
import json
from pptx import Presentation
from pptx.util import Pt

from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

def apply_korean_font(paragraph, font_name, font_size=None):
    """단락 내의 모든 런에 폰트를 강제 적용하는 헬퍼 함수"""
    for run in paragraph.runs:
        run.font.name = font_name
        
        # XML 수준에서 동아시아 글꼴 설정 (파워포인트 한글 폰트 강제 적용 로직)
        rPr = run.font._element
        rFonts = rPr.find(qn('a:rFonts'))
        if rFonts is None:
            rFonts = OxmlElement('a:rFonts')
            rPr.append(rFonts)
            
        rFonts.set('ea', font_name)
        rFonts.set('ascii', font_name)
        rFonts.set('hAnsi', font_name)
        rFonts.set('cs', font_name)
        
        if font_size:
            run.font.size = Pt(font_size)

def generate_ppt_from_json(json_data: dict, template_path: str = None, output_path: str = "output.pptx", font_name: str = "맑은 고딕"):
    """
    JSON 데이터를 기반으로 PPTX 파일을 생성 또는 조립하여 반환합니다.
    """
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation() # 기본 템플릿 사용
        
    slides_data = json_data.get('slides', [])
    topic_title = json_data.get('topic', '생성된 프레젠테이션')
    
    # 1. 메인 타이틀 슬라이드 (가장 첫 번째 레이아웃 사용)
    # 레이아웃이 1개라도 있으면 인덱스 0을 사용하고, 2개 이상일 때만 인덱스 1을 본문용으로 씁니다.
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    try:
        title = slide.shapes.title
        title.text = topic_title
        for paragraph in title.text_frame.paragraphs:
            apply_korean_font(paragraph, font_name, 40)
    except Exception:
        pass
    
    try:
        subtitle = slide.placeholders[1]
        subtitle.text = f"AI Generated Presentation\nGenerated at: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"
        for paragraph in subtitle.text_frame.paragraphs:
            apply_korean_font(paragraph, font_name, 18)
    except Exception:
        pass

    # 2. 본문 슬라이드 레이아웃 결정
    # 템플릿에 레이아웃이 1개뿐이면(표지 전용) 어쩔 수 없이 0번을 쓰고, 2개 이상이면 1번을 씁니다.
    bullet_slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        try:
            body_shape = shapes.placeholders[1]
        except KeyError:
            # Placeholder가 없으면 생성하지 않음 (엣지 케이스 방어)
            continue
            
        title_shape.text = slide_info.get("title", "제목 없음")
        # 제목 폰트 적용
        for paragraph in title_shape.text_frame.paragraphs:
            apply_korean_font(paragraph, font_name, 32)
            
        tf = body_shape.text_frame
        tf.clear() # 기본 플레이스홀더 텍스트 제거
        
        contents = slide_info.get("content", [])
        for i, bullet in enumerate(contents):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            apply_korean_font(p, font_name, 22)
            
    # 3. 메타데이터 (Core Properties) 은닉 삽입
    # 파일 속성(작성자, 태그, 설명 등)에 생성 상세 정보를 남깁니다.
    prs.core_properties.author = "PPT Agent (AI)"
    prs.core_properties.title = topic_title
    prs.core_properties.comments = json.dumps({
        "generated_at": datetime.datetime.now().isoformat(),
        "template_used": template_path if template_path else "default",
        "font_used": font_name
    }, ensure_ascii=False)
    prs.core_properties.keywords = f"AIGenerated, PPTAgent, {font_name}"
    
    prs.save(output_path)
    return output_path
